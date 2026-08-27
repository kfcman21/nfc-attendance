# -*- coding: utf-8 -*-
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_perfect_hackathon_pptx():
    prs = pptx.Presentation()
    # A4 세로 규격 (210mm x 297mm = 8.27in x 11.69in)
    prs.slide_width = Inches(8.27)
    prs.slide_height = Inches(11.69)
    blank_layout = prs.slide_layouts[6]

    # 공통 색상 팔레트
    C_DARK = RGBColor(40, 48, 60)        # 메인 다크 그레이
    C_YELLOW = RGBColor(254, 209, 86)    # 1/4번 노랑
    C_TEAL = RGBColor(85, 214, 190)      # 2번 청록
    C_PINK = RGBColor(252, 141, 158)     # 3번 핑크
    C_BLUE = RGBColor(0, 100, 224)       # 포인트 코발트 블루
    C_BORDER = RGBColor(205, 212, 220)   # 박스 외곽선
    C_BOX_BG = RGBColor(255, 255, 255)   # 박스 배경
    C_TEXT = RGBColor(25, 30, 36)        # 본문 텍스트
    C_MUTED = RGBColor(110, 120, 135)    # 보조 텍스트
    C_LIGHT_GRAY = RGBColor(245, 247, 250)

    def draw_base_template(slide, page_num, section_title, badge_color):
        # 1. 상단 다크 헤더 밴드
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(8.27), Inches(0.55))
        header.fill.solid()
        header.fill.fore_color.rgb = C_DARK
        header.line.color.rgb = C_DARK
        tf = header.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = "해커톤 팀 활동 기획서"
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # 2. 대주제 번호 배지
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(0.8), Inches(0.55), Inches(0.55))
        badge.fill.solid()
        badge.fill.fore_color.rgb = badge_color
        badge.line.fill.background()
        tf_b = badge.text_frame
        tf_b.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_b = tf_b.paragraphs[0]
        p_b.text = str(page_num)
        p_b.font.name = "Malgun Gothic"
        p_b.font.size = Pt(20)
        p_b.font.bold = True
        p_b.font.color.rgb = C_DARK
        p_b.alignment = PP_ALIGN.CENTER

        # 3. 대주제 타이틀
        tx_title = slide.shapes.add_textbox(Inches(1.3), Inches(0.8), Inches(4.0), Inches(0.55))
        tf_t = tx_title.text_frame
        tf_t.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_t = tf_t.paragraphs[0]
        p_t.text = section_title
        p_t.font.name = "Malgun Gothic"
        p_t.font.size = Pt(20)
        p_t.font.bold = True
        p_t.font.color.rgb = C_DARK

        # 4. 하단 푸터
        footer_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(11.15), Inches(6.97), Inches(0.01))
        footer_line.fill.solid()
        footer_line.fill.fore_color.rgb = RGBColor(220, 225, 230)
        footer_line.line.color.rgb = RGBColor(220, 225, 230)

        tx_foot_l = slide.shapes.add_textbox(Inches(0.65), Inches(11.2), Inches(5.0), Inches(0.35))
        tf_fl = tx_foot_l.text_frame
        p_fl = tf_fl.paragraphs[0]
        p_fl.text = "2026 기업 연계 정보교원 역량강화 프로젝트 · 해커톤 팀 기획서"
        p_fl.font.name = "Malgun Gothic"
        p_fl.font.size = Pt(9.5)
        p_fl.font.color.rgb = C_MUTED

        tx_foot_r = slide.shapes.add_textbox(Inches(6.62), Inches(11.2), Inches(1.0), Inches(0.35))
        tf_fr = tx_foot_r.text_frame
        p_fr = tf_fr.paragraphs[0]
        p_fr.text = f"{page_num} / 4"
        p_fr.font.name = "Malgun Gothic"
        p_fr.font.size = Pt(9.5)
        p_fr.font.bold = True
        p_fr.font.color.rgb = C_MUTED
        p_fr.alignment = PP_ALIGN.RIGHT

    def draw_info_bar(slide, is_first_page=False):
        top_y = Inches(1.48)
        height = Inches(0.7) if is_first_page else Inches(0.4)
        
        info_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), top_y, Inches(6.97), height)
        info_box.fill.solid()
        info_box.fill.fore_color.rgb = C_LIGHT_GRAY
        info_box.line.color.rgb = C_BORDER
        info_box.line.width = Pt(1)
        
        tf = info_box.text_frame
        tf.margin_left = Inches(0.15)
        tf.margin_top = Inches(0.08)
        tf.word_wrap = True
        
        if is_first_page:
            p1 = tf.paragraphs[0]
            p1.text = "분반: 초등 정보 융합 분반         모둠명: NFC 에듀태그 팀 (Smart EduLab)"
            p1.font.name = "Malgun Gothic"
            p1.font.size = Pt(10.5)
            p1.font.bold = True
            p1.font.color.rgb = C_DARK
            
            p2 = tf.add_paragraph()
            p2.text = "팀원: 박찬규(대표/개발 총괄), 팀원 A(교육과정/콘텐츠), 팀원 B(UI·UX 디자인/현장검증)"
            p2.font.name = "Malgun Gothic"
            p2.font.size = Pt(9.5)
            p2.font.color.rgb = C_BLUE
            p2.space_before = Pt(4)
        else:
            p1 = tf.paragraphs[0]
            p1.text = "모둠명: NFC 에듀태그 팀 (Smart EduLab)       |       대표: 박찬규"
            p1.font.name = "Malgun Gothic"
            p1.font.size = Pt(10)
            p1.font.bold = True
            p1.font.color.rgb = C_DARK

    def draw_section_header(slide, top_y, badge_text, sub_text):
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), top_y, Inches(1.4), Inches(0.32))
        badge.fill.solid()
        badge.fill.fore_color.rgb = C_DARK
        badge.line.fill.background()
        tf_b = badge.text_frame
        tf_b.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_b = tf_b.paragraphs[0]
        p_b.text = badge_text
        p_b.font.name = "Malgun Gothic"
        p_b.font.size = Pt(10.5)
        p_b.font.bold = True
        p_b.font.color.rgb = RGBColor(255, 255, 255)
        p_b.alignment = PP_ALIGN.CENTER

        tx_sub = slide.shapes.add_textbox(Inches(2.15), top_y - Inches(0.02), Inches(5.4), Inches(0.35))
        tf_s = tx_sub.text_frame
        tf_s.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_s = tf_s.paragraphs[0]
        p_s.text = sub_text
        p_s.font.name = "Malgun Gothic"
        p_s.font.size = Pt(10)
        p_s.font.italic = True
        p_s.font.color.rgb = C_MUTED

    def draw_content_box(slide, top_y, height, text_list, font_size=Pt(9.5), line_space=Pt(4)):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), top_y, Inches(6.97), height)
        box.fill.solid()
        box.fill.fore_color.rgb = C_BOX_BG
        box.line.color.rgb = C_BORDER
        box.line.width = Pt(1)
        
        tf = box.text_frame
        tf.margin_left = Inches(0.18)
        tf.margin_right = Inches(0.18)
        tf.margin_top = Inches(0.12)
        tf.margin_bottom = Inches(0.12)
        tf.word_wrap = True
        
        for i, item in enumerate(text_list):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.font.name = "Malgun Gothic"
            p.font.size = font_size
            p.font.color.rgb = C_TEXT
            if i > 0:
                p.space_before = line_space
        return box


    # =========================================================================
    # SLIDE 1: 문제 찾기
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    draw_base_template(s1, 1, "문제 찾기", C_YELLOW)
    draw_info_bar(s1, is_first_page=True)

    # 1. 문제 정의
    draw_section_header(s1, Inches(2.32), "1. 문제 정의", "누구의 어떤 페인포인트인지")
    draw_content_box(s1, Inches(2.70), Inches(2.45), [
        "• 교사의 과중한 학급 및 과학실 관리 행정력 소모:",
        "  - 매일 아침 학생 출결 수동 체크, 지각생 파악, 학생들의 정서·마음 상태를 일일이 파악하기 어려움.",
        "• 지능형 과학실 탐구 이력 및 안전 관리 부재:",
        "  - 과학실 스테이션(MBL, 현미경 등)별 학생 탐구 이력 관리의 어려움 및 고가 교구·화학 시약(MSDS) 대여/반납과 실험 안전 수칙 안내의 형식화.",
        "• 분절된 학급 관리 도구의 비효율:",
        "  - 출석부, 독서 통장, 체육(PAPS) 측정, 학급 화폐가 제각각 분리되어 교사의 수업 준비 및 관리 피로도 가중.",
        "• 교육용 AI의 무비판적 완전 자동화 위험:",
        "  - AI가 생성한 피드백이 교사의 검토 없이 학생에게 그대로 전달되어, 교사의 주체성과 학생의 비판적 사고가 함께 약화됨."
    ], font_size=Pt(9.0), line_space=Pt(2))

    # 2. 지금 방식
    draw_section_header(s1, Inches(5.30), "2. 지금 방식", "지금은 어떻게 하고 있는지 · 불편한 점")
    draw_content_box(s1, Inches(5.68), Inches(2.45), [
        "• 출결 및 감정: 종이 출석부나 나이스(NEIS) 수동 입력에 의존하여 학생의 정서 변화나 정시 등교 추이를 직관적으로 파악하기 어려움.",
        "• 과학실 탐구: 종이 활동지 도장 날인 및 수기 대여 대장에 의존하여 교구 분실 위험이 높고 탐구 데이터 누적 불가능.",
        "• 체육 및 독서: 초시계를 들고 수동으로 셔틀런 랩타임을 측정하거나 종이에 적어 체계적인 성장 분석이 어려움.",
        "• 불편한 점: 클라우드 웹앱은 교실 인터넷 불안정 시 접속이 끊기며, 학생 개인정보(이름, 사진) 유출 우려가 상존함."
    ], font_size=Pt(9.3), line_space=Pt(3))

    # 3. 사용자
    draw_section_header(s1, Inches(8.28), "3. 사용자", "누가 · 어떤 상황에서 쓰는지")
    draw_content_box(s1, Inches(8.66), Inches(2.25), [
        "• 주 사용자 1 (초·중등 교사):",
        "  - 아침 등교 맞이(출석/감정 AI 케어), 지능형 과학실 스테이션 탐구 수업 및 교구/시약 안전 통제, 체육 PAPS 측정, 학급 문고 및 1인 1역할 학급 화폐(지구 포인트)를 통합 정산할 때.",
        "• 주 사용자 2 (초등학생):",
        "  - 등교 시, 과학실 탐구 코너 이동 시, 교구 대여 시 본인의 NFC 학생 카드를 리더기에 '톡' 태그할 때."
    ], font_size=Pt(9.3), line_space=Pt(3))


    # =========================================================================
    # SLIDE 2: 기능과 자료
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    draw_base_template(s2, 2, "기능과 자료", C_TEAL)
    draw_info_bar(s2, is_first_page=False)

    # 4. 필요 기능
    draw_section_header(s2, Inches(2.00), "4. 필요 기능", "핵심 기능 위주 · 한 줄에 하나")
    draw_content_box(s2, Inches(2.38), Inches(4.00), [
        "✔  NFC 원터치 정시 출석 및 점수제 : 마감시각 이전 태그 시 만점 부여 및 실시간 대시보드 시각화",
        "✔  Upstage Solar AI 감정 출석부 : 출석 후 기분 선택 시 초거대 AI의 맞춤형 긍정 응원 메시지 실시간 생성",
        "✔  스마트 과학 탐구 패스포트 : 과학실 스테이션별 NFC 태깅 미션 완주 스탬프 랠리 자동화",
        "✔  스마트 교구 & 시약 안전 지킴이 : NFC 교구 원터치 대여/반납 및 MSDS 실험 안전 수칙 실시간 팝업 안내",
        "✔  SciBit 마이크로비트 MBL 연계 : micro:bit v2 센서(온도·조도·소음·가속도) 실시간 수집 및 기록 (kfcman.link/scibit)",
        "✔  SciBot 햄스터 피지컬 AI 탐사 연계 : 햄스터 로봇 온실 순찰 & 자율주행 탐사 미션 스탬프 발급 (kfcman.link/scibot)",
        "✔  스마트 PAPS 셔틀런 & 서킷 : 반환점 태깅으로 왕복달리기 자동 기록 및 Solar AI 체력 코칭 리포트",
        "✔  손바닥 도서관 & 지구 포인트 상점 : 학급문고 원터치 대여 및 1인 1역할 보상 학급 화폐 경제 실습",
        "✔  인간 협업형(HITL) AI 교육 스튜디오 : AI-SPARC 수업 설계기 + 교사 Level 3 검토·승인 센터 + AI 감사 추적 활동지 제작기",
        "✔  🖥️ 바탕화면 달력 연동 : 설정 탭에서 데스크톱 달력 앱 설치 파일(exe) 원클릭 연동 실행",
        "✔  공공데이터 & 100% 로컬 보안 : 기상청·에어코리아·NEIS 실시간 연동 및 완전 오프라인 SQLite 암호화 보관"
    ], font_size=Pt(8.2), line_space=Pt(2))

    # 5. 입력 자료
    draw_section_header(s2, Inches(6.52), "5. 입력 자료", "넣을 자료 · 데이터 · 예시 문서")
    draw_content_box(s2, Inches(6.90), Inches(1.80), [
        "• 학생 및 교구 데이터: 학생 번호, 카드 UID, 과학실 스테이션 정보, 교구/시약 보관 위치 및 MSDS 안전 수칙.",
        "• 센서 및 공공데이터: micro:bit MBL 센서값, 기상청 단기예보(기온/날씨), 에어코리아 미세먼지(PM10/PM2.5), NEIS 급식.",
        "• 보안 API 키: Upstage Solar API Key (AES-256-GCM 암호화 보관)."
    ], font_size=Pt(9.3), line_space=Pt(3))

    # 6. 만들 방식
    draw_section_header(s2, Inches(8.85), "6. 만들 방식", "택1 또는 혼합")
    draw_content_box(s2, Inches(9.23), Inches(0.65), [
        "☐ 에이전트          ☐ 바이브코딩          ☑ 혼합 (Hybrid) — [Electron/Node 바이브코딩 + Upstage Solar AI 에이전트]"
    ], font_size=Pt(9.5))

    # 7. 사용 도구
    draw_section_header(s2, Inches(10.02), "7. 사용 도구", "쓸 도구 이름")
    draw_content_box(s2, Inches(10.40), Inches(0.60), [
        "Node.js, Express, Electron, electron-builder, SQLite3, SerialPort, HTML5/CSS3(진한 굴림체), Upstage Solar LLM API"
    ], font_size=Pt(9.5))


    # =========================================================================
    # SLIDE 3: 동작 흐름
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    draw_base_template(s3, 3, "동작 흐름", C_PINK)
    draw_info_bar(s3, is_first_page=False)

    draw_section_header(s3, Inches(2.00), "8. 동작 흐름", "단계별 내용 · 글 또는 그림")

    # 3단 컬럼 박스 생성
    col_w = Inches(2.05)
    col_h = Inches(4.30)
    col_y = Inches(2.40)

    # 컬럼 1: 입력
    b1 = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), col_y, col_w, col_h)
    b1.fill.solid()
    b1.fill.fore_color.rgb = RGBColor(255, 250, 235)
    b1.line.color.rgb = C_YELLOW
    b1.line.width = Pt(1.5)
    tf1 = b1.text_frame
    tf1.margin_left = Inches(0.12)
    tf1.margin_top = Inches(0.12)
    tf1.word_wrap = True
    p1_h = tf1.paragraphs[0]
    p1_h.text = "1. 입 력 (Input)"
    p1_h.font.name = "Malgun Gothic"
    p1_h.font.size = Pt(11)
    p1_h.font.bold = True
    p1_h.font.color.rgb = RGBColor(160, 110, 0)
    p1_h.alignment = PP_ALIGN.CENTER
    
    items1 = [
        "",
        "• NFC 카드 태그",
        "  - 학생 카드 (학생 식별)",
        "  - 스테이션 카드 (완주)",
        "  - 교구 카드 (대여/반납)",
        "  - 책 카드 (도서 대여)",
        "",
        "• micro:bit MBL 센서",
        "  - 온도, 조도, 소음, 가속도",
        "",
        "• 햄스터 로봇 미션",
        "  - 온실 순찰 / 자율주행",
        "",
        "• 바탕화면 달력 실행 클릭"
    ]
    for it in items1:
        p = tf1.add_paragraph()
        p.text = it
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(8.2)
        p.font.color.rgb = C_TEXT

    # 화살표 1
    arr1 = s3.shapes.add_textbox(Inches(2.75), Inches(4.2), Inches(0.35), Inches(0.5))
    arr1.text_frame.paragraphs[0].text = "➔"
    arr1.text_frame.paragraphs[0].font.size = Pt(18)
    arr1.text_frame.paragraphs[0].font.bold = True
    arr1.text_frame.paragraphs[0].font.color.rgb = C_MUTED

    # 컬럼 2: 처리
    b2 = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.11), col_y, col_w, col_h)
    b2.fill.solid()
    b2.fill.fore_color.rgb = RGBColor(235, 252, 248)
    b2.line.color.rgb = C_TEAL
    b2.line.width = Pt(1.5)
    tf2 = b2.text_frame
    tf2.margin_left = Inches(0.12)
    tf2.margin_top = Inches(0.12)
    tf2.word_wrap = True
    p2_h = tf2.paragraphs[0]
    p2_h.text = "2. 처 리 (Process)"
    p2_h.font.name = "Malgun Gothic"
    p2_h.font.size = Pt(11)
    p2_h.font.bold = True
    p2_h.font.color.rgb = RGBColor(0, 130, 110)
    p2_h.alignment = PP_ALIGN.CENTER
    
    items2 = [
        "",
        "• 리더기 모드별 분기",
        "  - 출석/과학실/도서/체육",
        "",
        "• 로컬 SQLite 트랜잭션",
        "  - 출결, 스탬프, 대여 기록",
        "",
        "• Upstage Solar AI 분석",
        "  - 감정 체크인 응원 생성",
        "  - 셔틀런 맞춤 코칭",
        "",
        "• HITL Level 3 교사 검토",
        "  - AI 초안 비판·수정·승인",
        "",
        "• 실시간 SSE 브로드캐스트",
        "",
        "• 달력 exe spawn 프로세스 구동"
    ]
    for it in items2:
        p = tf2.add_paragraph()
        p.text = it
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(8.2)
        p.font.color.rgb = C_TEXT

    # 화살표 2
    arr2 = s3.shapes.add_textbox(Inches(5.21), Inches(4.2), Inches(0.35), Inches(0.5))
    arr2.text_frame.paragraphs[0].text = "➔"
    arr2.text_frame.paragraphs[0].font.size = Pt(18)
    arr2.text_frame.paragraphs[0].font.bold = True
    arr2.text_frame.paragraphs[0].font.color.rgb = C_MUTED

    # 컬럼 3: 결과
    b3 = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.57), col_y, col_w, col_h)
    b3.fill.solid()
    b3.fill.fore_color.rgb = RGBColor(255, 242, 245)
    b3.line.color.rgb = C_PINK
    b3.line.width = Pt(1.5)
    tf3 = b3.text_frame
    tf3.margin_left = Inches(0.12)
    tf3.margin_top = Inches(0.12)
    tf3.word_wrap = True
    p3_h = tf3.paragraphs[0]
    p3_h.text = "3. 결 과 (Output)"
    p3_h.font.name = "Malgun Gothic"
    p3_h.font.size = Pt(11)
    p3_h.font.bold = True
    p3_h.font.color.rgb = RGBColor(180, 50, 80)
    p3_h.alignment = PP_ALIGN.CENTER
    
    items3 = [
        "",
        "• 대형 모니터 실시간 안내",
        "  - 정시 출석 완료 & 점수",
        "  - Solar AI 맞춤 격려 팝업",
        "  - 패스포트 스탬프 발급",
        "  - 교구 대여 & MSDS 표출",
        "",
        "  - 교사 승인 피드백만 배포",
        "",
        "• AI 감사 추적 활동지 출력",
        "  - 탐구→질문→태깅→교정",
        "",
        "• 원클릭 엑셀(CSV) 저장",
        "  - 나이스/생활기록부 반영",
        "  - 100% 로컬 데이터 보존",
        "",
        "• 바탕화면 달력 연동화면 표출"
    ]
    for it in items3:
        p = tf3.add_paragraph()
        p.text = it
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(8.2)
        p.font.color.rgb = C_TEXT

    # 9. 멘토 질문
    draw_section_header(s3, Inches(6.90), "9. 멘토 질문", "막힌 것 · 멘토에게 물어볼 것")
    draw_content_box(s3, Inches(7.28), Inches(3.60), [
        "1. 하드웨어 호환성 및 안정성 검증:",
        "   - 학교 컴퓨터실마다 OS 환경(Win 10/11)과 시리얼 COM 포트 번호가 다른데, 시리얼 포트 자동 탐지(Auto-detect) 및 USB HID 키보드 모드 겸용 구조가 실제 현장에서 안정적으로 구동되는지 점검받고 싶습니다.",
        "",
        "2. 피지컬 컴퓨팅 교구 다이렉트 연동 정책:",
        "   - 지능형 과학실의 마이크로비트(SciBit) 및 햄스터 로봇(SciBot) 데이터를 브라우저의 Web Serial/Web Bluetooth API로 다이렉트 연동할 때의 보안 가이드라인 및 모범 사례에 대한 자문을 구합니다.",
        "",
        "3. 오프라인 AI 폴백(Fallback) 아키텍처:",
        "   - 학교 교실 내 인터넷이 일시적으로 끊겼을 때 Upstage Solar AI 호출 실패를 대비하여, 오프라인 룰베이스(Rule-based) 감정 및 체육 피드백으로 부드럽게 전환되는 Fallback 설계가 적절한지 평가받고 싶습니다."
    ], font_size=Pt(9.0), line_space=Pt(2))


    # =========================================================================
    # SLIDE 4: 목표 정하기
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    draw_base_template(s4, 4, "목표 정하기", C_YELLOW)
    draw_info_bar(s4, is_first_page=False)

    # 10. 결과물 이름
    draw_section_header(s4, Inches(2.00), "10. 결과물 이름", "가칭 · 나중에 바꿔도 됨")
    box10 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(2.38), Inches(6.97), Inches(0.65))
    box10.fill.solid()
    box10.fill.fore_color.rgb = RGBColor(238, 246, 255)
    box10.line.color.rgb = C_BLUE
    box10.line.width = Pt(1.5)
    tf10 = box10.text_frame
    tf10.vertical_anchor = MSO_ANCHOR.MIDDLE
    p10 = tf10.paragraphs[0]
    p10.text = "NFC 에듀태그 (NFC EduTag) - 스마트 학급 & 지능형 과학실 올인원 솔루션"
    p10.font.name = "Malgun Gothic"
    p10.font.size = Pt(12)
    p10.font.bold = True
    p10.font.color.rgb = C_BLUE
    p10.alignment = PP_ALIGN.CENTER

    # 11. 한 줄 소개
    draw_section_header(s4, Inches(3.20), "11. 한 줄 소개", "만들려는 것 한 문장")
    box11 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(3.58), Inches(6.97), Inches(0.75))
    box11.fill.solid()
    box11.fill.fore_color.rgb = C_LIGHT_GRAY
    box11.line.color.rgb = C_BORDER
    tf11 = box11.text_frame
    tf11.vertical_anchor = MSO_ANCHOR.MIDDLE
    p11 = tf11.paragraphs[0]
    p11.text = '"카드 한 장으로 학생 출결·감정 케어부터 지능형 과학실 MBL 탐구·로봇 피지컬 AI까지 원터치로 해결하는 미래형 스마트 스쿨 올인원 솔루션"'
    p11.font.name = "Malgun Gothic"
    p11.font.size = Pt(10)
    p11.font.bold = True
    p11.font.color.rgb = C_DARK
    p11.alignment = PP_ALIGN.CENTER

    # 12. 역할 분담
    draw_section_header(s4, Inches(4.50), "12. 역할 분담", "팀원별 담당")
    
    # 역할 분담 표 생성
    tbl_shape = s4.shapes.add_table(4, 2, Inches(0.65), Inches(4.88), Inches(6.97), Inches(2.40))
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(1.8)
    tbl.columns[1].width = Inches(5.17)

    headers = ["이 름", "담당 업무 및 역할"]
    for c_idx, h_text in enumerate(headers):
        cell = tbl.cell(0, c_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_DARK
        p = cell.text_frame.paragraphs[0]
        p.text = h_text
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    rows_data = [
        ("박찬규\n(대표 / 개발 총괄)", "• 전체 시스템 아키텍처 설계, Electron 데스크톱 앱 및 Node/SQLite DB 구축\n• 지능형 과학실(SciBit/SciBot) 연계 및 Upstage Solar AI 엔진 통합 개발\n• 2022 개정 교육과정 연계 기획 및 최종 독립형 패키징 빌드"),
        ("팀원 A\n(교육과정 / 콘텐츠)", "• 2022 개정 교육과정 연계 교과별 탐구 미션 및 교수·학습 활동지 개발\n• HITL AI-SPARC 수업 설계안 및 4단계 AI 감사 추적 학생 활동지 집필\n• 과학 교구 및 화학 시약 MSDS 물질안전보건자료 텍스트 DB 정리"),
        ("팀원 B\n(UI·UX 디자인 / 검증)", "• 공식 3D 인포그래픽 포스터 디자인 및 전자칠판 최적화 진한 굴림체 UI 튜닝\n• 학교 현장 시연 테스트 및 기능별 트러블슈팅/사용자 피드백 수집")
    ]

    for r_idx, (name, role) in enumerate(rows_data, start=1):
        cell_n = tbl.cell(r_idx, 0)
        cell_n.fill.solid()
        cell_n.fill.fore_color.rgb = C_LIGHT_GRAY if r_idx % 2 == 1 else RGBColor(255, 255, 255)
        p_n = cell_n.text_frame.paragraphs[0]
        p_n.text = name
        p_n.font.name = "Malgun Gothic"
        p_n.font.size = Pt(9)
        p_n.font.bold = True
        p_n.font.color.rgb = C_BLUE
        p_n.alignment = PP_ALIGN.CENTER

        cell_r = tbl.cell(r_idx, 1)
        cell_r.fill.solid()
        cell_r.fill.fore_color.rgb = C_LIGHT_GRAY if r_idx % 2 == 1 else RGBColor(255, 255, 255)
        cell_r.text_frame.margin_left = Inches(0.1)
        cell_r.text_frame.margin_right = Inches(0.1)
        p_r = cell_r.text_frame.paragraphs[0]
        p_r.text = role
        p_r.font.name = "Malgun Gothic"
        p_r.font.size = Pt(8.5)
        p_r.font.color.rgb = C_TEXT

    # 13. 완성 기준
    draw_section_header(s4, Inches(7.45), "13. 완성 기준", "이것이 되면 완성")
    draw_content_box(s4, Inches(7.83), Inches(3.25), [
        "1. [필수] 완전 독립형 배포: 인터넷이나 추가 설치 없이 학교 PC에서 NFC에듀태그.exe 더블클릭만으로 100% 정상 구동.",
        "2. [필수] NFC 올인원 다중 모드: 학생 카드 태그 시 출석·감정 체크, 과학실 패스포트 스탬프 발급, 교구 대여/반납이 1초 내에 오류 없이 처리.",
        "3. [필수] 지능형 과학실 융합: kfcman.link/scibit micro:bit 센서 수집 및 kfcman.link/scibot 햄스터 로봇 탐사 미션이 패스포트에 정상 연동.",
        "4. [필수] 초거대 AI 및 로컬 보안: Upstage Solar AI 피드백이 실시간 생성되고, API 키와 학생 정보가 AES-256 로컬 암호화로 안전하게 보관.",
        "5. [필수] 교사용 데이터 관리: 모든 활동 기록이 클릭 한 번으로 엑셀(CSV)로 즉시 다운로드되고 전체 백업/복원이 완벽 지원.",
        "6. [필수] 인간 협업형(HITL) 검증: AI 피드백이 교사의 Level 3 검토·승인을 거친 뒤에만 학생에게 배포되며, AI-SPARC 수업 설계안과 4단계 AI 감사 추적 활동지가 정상 생성.",
        "7. [선택] 바탕화면 달력 연동: 데스크톱 달력 설치 파일(exe) 경로 설정 시 원클릭 프로세스 실행 기능 정상 작동."
    ], font_size=Pt(8.2), line_space=Pt(1.5))

    out_path1 = r"C:\Users\박찬규\Desktop\(완성본)해커톤_팀_활동_기획서_NFC에듀태그.pptx"
    out_path2 = r"C:\Users\박찬규\Desktop\Project\nfc\docs\(완성본)해커톤_팀_활동_기획서_NFC에듀태그.pptx"
    out_path3 = r"C:\Users\박찬규\Desktop\Project\nfc\(완성본)해커톤_팀_활동_기획서_NFC에듀태그.pptx"

    prs.save(out_path1)
    prs.save(out_path2)
    prs.save(out_path3)
    print("성공: 글자 맞춤 최적화 폼이 반영된 해커톤 기획서 PPTX가 생성되었습니다!")

    # comtypes를 사용한 PPTX -> PDF 자동 변환
    try:
        import comtypes.client
        import os
        print("PDF 자동 변환을 시작합니다...")
        
        # PPTX 파일 존재 확인 후 변환 진행
        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
        powerpoint.Visible = True  # COM 오토메이션 안정성을 위해 창 활성화
        
        for pptx_path in [out_path1, out_path2, out_path3]:
            if os.path.exists(pptx_path):
                pdf_path = pptx_path.replace(".pptx", ".pdf")
                abs_pptx = os.path.abspath(pptx_path)
                abs_pdf = os.path.abspath(pdf_path)
                
                print(f"변환 진행 중: {abs_pptx} -> {abs_pdf}")
                deck = powerpoint.Presentations.Open(abs_pptx, WithWindow=False)
                deck.SaveAs(abs_pdf, 32)  # ppSaveAsPDF = 32
                deck.Close()
                print(f"변환 성공: {pdf_path}")
                
        powerpoint.Quit()
        print("성공: 모든 PPTX 파일이 PDF로 성공적으로 자동 변환되었습니다!")
    except Exception as e:
        print(f"경고: PDF 자동 변환 중 에러 발생 (수동 변환 필요할 수 있음): {e}")

if __name__ == "__main__":
    create_perfect_hackathon_pptx()
