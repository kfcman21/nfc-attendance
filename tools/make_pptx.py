# -*- coding: utf-8 -*-
import sys
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def fill_hackathon_pptx():
    src_path = r"C:\Users\박찬규\Desktop\(워크시트)해커톤_기획서.pptx"
    out_path1 = r"C:\Users\박찬규\Desktop\(완성본)해커톤_팀_활동_기획서_NFC에듀태그.pptx"
    out_path2 = r"C:\Users\박찬규\Desktop\Project\nfc\docs\(완성본)해커톤_팀_활동_기획서_NFC에듀태그.pptx"
    out_path3 = r"C:\Users\박찬규\Desktop\Project\nfc\(완성본)해커톤_팀_활동_기획서_NFC에듀태그.pptx"

    prs = pptx.Presentation(src_path)

    def add_textbox(slide, left, top, width, height, text, font_size=Pt(11), bold=False, color=RGBColor(30, 30, 30), align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.08)
        tf.margin_bottom = Inches(0.08)
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = "Malgun Gothic"
        p.font.size = font_size
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = align
        return txBox

    # ==================== SLIDE 1: 문제 찾기 ====================
    s1 = prs.slides[0]
    # 분반, 모둠명, 팀원 정보
    add_textbox(s1, Inches(1.2), Inches(1.58), Inches(1.5), Inches(0.3), "초등 정보 융합 분반", Pt(10), True, RGBColor(0, 80, 180))
    add_textbox(s1, Inches(4.3), Inches(1.58), Inches(2.2), Inches(0.3), "NFC 에듀태그 팀", Pt(10), True, RGBColor(0, 80, 180))
    add_textbox(s1, Inches(1.2), Inches(1.92), Inches(5.8), Inches(0.3), "박찬규(대표/개발), 팀원A(교육과정), 팀원B(디자인/검증)", Pt(9.5), True, RGBColor(0, 80, 180))

    # 1. 문제 정의 (박스 내부)
    text_p1_1 = (
        "• 교사의 과중한 학급/과학실 관리 행정력 소모: 매일 아침 학생 출결 체크, 지각생 파악, 학생 마음 상태 파악의 번거로움.\n"
        "• 지능형 과학실 탐구 이력 및 안전 관리 부재: 과학실 스테이션(MBL, 현미경 등)별 탐구 이력 관리의 어려움 및 고가 교구/화학 시약(MSDS) 대여/반납과 안전 수칙 안내의 형식화.\n"
        "• 분절된 학급 도구: 출석부, 독서 통장, 체육(PAPS) 측정, 학급 화폐가 제각각 분리되어 교사 업무 피로도 누적."
    )
    add_textbox(s1, Inches(0.65), Inches(2.65), Inches(6.2), Inches(1.2), text_p1_1, Pt(9.5), False, RGBColor(20, 20, 20))

    # 2. 지금 방식
    text_p1_2 = (
        "• 출결/감정: 종이 출석부나 나이스(NEIS) 수동 입력에 의존하여 학생의 정서 변화나 정시 등교 추이를 직관적으로 파악하기 어려움.\n"
        "• 과학실 탐구: 종이 활동지 도장 날인 및 수기 대여 대장에 의존하여 교구 분실 위험이 높고 탐구 데이터 누적 불가능.\n"
        "• 체육/독서: 초시계를 들고 수동으로 셔틀런 랩타임을 측정하거나 종이에 적어 체계적인 성장 분석이 어려움.\n"
        "• 불편한 점: 클라우드 웹앱은 교실 인터넷 불안정 시 접속이 끊기며, 학생 개인정보(이름, 사진) 유출 우려 상존."
    )
    add_textbox(s1, Inches(0.65), Inches(4.35), Inches(6.2), Inches(1.2), text_p1_2, Pt(9.5), False, RGBColor(20, 20, 20))

    # 3. 사용자
    text_p1_3 = (
        "• 초·중등 교사: 아침 등교 맞이(출석/감정 AI 케어), 지능형 과학실 스테이션 탐구 수업 및 교구/시약 안전 통제, 체육 PAPS 측정, 학급 문고 및 1인 1역할 학급 화폐(지구 포인트)를 통합 정산할 때.\n"
        "• 초등학생: 등교 시, 과학실 탐구 코너 이동 시, 교구 대여 시 본인의 NFC 학생 카드를 리더기에 '톡' 태그할 때."
    )
    add_textbox(s1, Inches(0.65), Inches(6.05), Inches(6.2), Inches(1.1), text_p1_3, Pt(9.5), False, RGBColor(20, 20, 20))


    # ==================== SLIDE 2: 기능과 자료 ====================
    s2 = prs.slides[1]
    add_textbox(s2, Inches(4.8), Inches(0.85), Inches(2.2), Inches(0.3), "NFC 에듀태그 팀", Pt(10), True, RGBColor(0, 80, 180))

    # 4. 필요 기능 (체크리스트)
    features = [
        "NFC 원터치 정시 출석 및 점수제: 마감시각 이전 태그 시 만점 부여 및 실시간 대시보드 시각화",
        "Upstage Solar AI 감정 출석부: 출석 후 기분 선택 시 초거대 AI의 맞춤형 긍정 응원 메시지 생성",
        "스마트 과학 탐구 패스포트: 과학실 스테이션별 NFC 태깅 미션 완주 스탬프 랠리 자동화",
        "스마트 교구 & 시약 안전 지킴이: NFC 교구 원터치 대여/반납 및 MSDS 실험 안전 수칙 실시간 팝업",
        "SciBit 마이크로비트 MBL 연계: micro:bit v2 센서(온도·조도·소음·가속도) 실시간 수집 및 기록 (kfcman.link/scibit)",
        "SciBot 햄스터 피지컬 AI 탐사 연계: 햄스터 로봇 온실 순찰 & 자율주행 탐사 미션 스탬프 발급 (kfcman.link/scibot)",
        "스마트 PAPS 셔틀런 & 서킷: 반환점 태깅으로 왕복달리기 자동 기록 및 Solar AI 체력 코칭 리포트",
        "손바닥 도서관 & 지구 포인트 상점: 학급문고 원터치 대여 및 1인 1역할 보상 학급 화폐 경제 실습"
    ]
    y_start = 1.62
    for feat in features:
        add_textbox(s2, Inches(0.88), Inches(y_start), Inches(6.0), Inches(0.24), feat, Pt(9), True, RGBColor(20, 20, 20))
        # 체크 표시
        add_textbox(s2, Inches(0.66), Inches(y_start - 0.02), Inches(0.25), Inches(0.24), "✔", Pt(10), True, RGBColor(0, 120, 215))
        y_start += 0.25

    # 5. 입력 자료
    text_p2_2 = (
        "• 학생 및 교구 데이터: 학생 번호, 카드 UID, 과학실 스테이션 정보, 교구/시약 보관 위치 및 MSDS 안전 수칙.\n"
        "• 센서 및 공공데이터: micro:bit MBL 센서값, 기상청 단기예보(기온/날씨), 에어코리아 미세먼지(PM10/PM2.5), NEIS 급식.\n"
        "• 보안 API 키: Upstage Solar API Key (AES-256-GCM 암호화 보관)."
    )
    add_textbox(s2, Inches(0.65), Inches(4.05), Inches(6.2), Inches(0.95), text_p2_2, Pt(9), False, RGBColor(20, 20, 20))

    # 6. 만들 방식 (혼합에 체크)
    add_textbox(s2, Inches(3.3), Inches(5.42), Inches(0.25), Inches(0.25), "✔", Pt(11), True, RGBColor(0, 120, 215))
    add_textbox(s2, Inches(4.0), Inches(5.42), Inches(3.0), Inches(0.3), "혼합 (Electron 바이브코딩 + Solar AI 에이전트)", Pt(9.5), True, RGBColor(0, 80, 180))

    # 7. 사용 도구
    text_p2_4 = "Node.js, Express, Electron, electron-builder, SQLite3, SerialPort, HTML5/CSS3(진한 굴림체), Upstage Solar LLM API"
    add_textbox(s2, Inches(0.65), Inches(6.28), Inches(6.2), Inches(0.55), text_p2_4, Pt(9.5), True, RGBColor(20, 20, 20))


    # ==================== SLIDE 3: 동작 흐름 ====================
    s3 = prs.slides[2]
    add_textbox(s3, Inches(4.8), Inches(0.85), Inches(2.2), Inches(0.3), "NFC 에듀태그 팀", Pt(10), True, RGBColor(0, 80, 180))

    # 8. 동작 흐름 (3단 컬럼)
    in_text = "• NFC 카드 태그\n  - 학생 카드 (식별)\n  - 스테이션 카드 (완료)\n  - 교구 카드 (대여/반납)\n  - 책 카드 (도서 대여)\n• micro:bit MBL 센서\n  - 온도, 조도, 소음, 가속도\n• 햄스터 로봇 탐사 출동"
    pro_text = "• 리더기 모드별 분기\n  - 출석/과학실/도서/체육\n• 로컬 SQLite 트랜잭션\n  - 출결, 스탬프, 대여 기록\n• Upstage Solar AI 분석\n  - 감정 체크인 & 체력 피드백\n• SSE 실시간 브로드캐스트"
    out_text = "• 대형 모니터 실시간 안내\n  - 출석 완료 & 획득 점수\n  - Solar AI 따뜻한 격려\n  - 패스포트 스탬프 (💮)\n  - 교구 대여 현황 & MSDS\n• 원클릭 엑셀(CSV) 저장\n  - 나이스/생기부 반영"

    add_textbox(s3, Inches(0.65), Inches(2.15), Inches(1.85), Inches(2.0), in_text, Pt(8.5), False, RGBColor(20, 20, 20))
    add_textbox(s3, Inches(2.80), Inches(2.15), Inches(1.85), Inches(2.0), pro_text, Pt(8.5), False, RGBColor(20, 20, 20))
    add_textbox(s3, Inches(4.95), Inches(2.15), Inches(1.95), Inches(2.0), out_text, Pt(8.5), False, RGBColor(20, 20, 20))

    # 9. 멘토 질문
    text_p3_2 = (
        "1. 하드웨어 호환성: 학교 컴퓨터실마다 OS(Win 10/11)와 COM 포트 번호가 다른데, 시리얼 포트 자동 감지(Auto-detect) 및 USB HID 키보드 모드 겸용 구조가 안정적인지 점검받고 싶습니다.\n"
        "2. 피지컬 교구 다이렉트 연동: 지능형 과학실의 마이크로비트(SciBit) 및 햄스터 로봇(SciBot) 데이터를 브라우저의 Web Serial/Bluetooth API로 연동할 때의 보안 가이드라인이 궁금합니다.\n"
        "3. 오프라인 AI 폴백(Fallback): 교실 내 인터넷이 일시적으로 끊겼을 때 Upstage Solar AI 호출 실패를 대비한 오프라인 룰베이스 피드백 전환 설계가 적절한지 자문을 구합니다."
    )
    add_textbox(s3, Inches(0.65), Inches(5.25), Inches(6.2), Inches(1.6), text_p3_2, Pt(9.5), False, RGBColor(20, 20, 20))


    # ==================== SLIDE 4: 목표 정하기 ====================
    s4 = prs.slides[3]
    add_textbox(s4, Inches(4.8), Inches(0.85), Inches(2.2), Inches(0.3), "NFC 에듀태그 팀", Pt(10), True, RGBColor(0, 80, 180))

    # 10. 결과물 이름
    add_textbox(s4, Inches(0.65), Inches(1.62), Inches(6.2), Inches(0.35), "NFC 에듀태그 (NFC EduTag) - 스마트 학급 & 지능형 과학실 올인원 솔루션", Pt(12), True, RGBColor(0, 80, 180))

    # 11. 한 줄 소개
    add_textbox(s4, Inches(0.65), Inches(2.55), Inches(6.2), Inches(0.55), '"카드 한 장으로 학생 출결·감정 케어부터 지능형 과학실 MBL 탐구·로봇 피지컬 AI까지 원터치로 해결하는 미래형 스마트 스쿨 올인원 솔루션"', Pt(10.5), True, RGBColor(20, 20, 20))

    # 12. 역할 분담
    add_textbox(s4, Inches(0.65), Inches(3.72), Inches(1.5), Inches(0.4), "박찬규 (대표/개발)", Pt(9.5), True, RGBColor(0, 80, 180))
    add_textbox(s4, Inches(2.2), Inches(3.68), Inches(4.6), Inches(0.5), "전체 아키텍처 설계, Electron 데스크톱 앱 및 Node 서버/SQLite DB 구축, 지능형 과학실(SciBit/SciBot) 연계, Upstage Solar AI 엔진 통합 및 패키징 빌드", Pt(8.5), False, RGBColor(20, 20, 20))

    add_textbox(s4, Inches(0.65), Inches(4.30), Inches(1.5), Inches(0.4), "팀원 A (교육과정)", Pt(9.5), True, RGBColor(0, 80, 180))
    add_textbox(s4, Inches(2.2), Inches(4.26), Inches(4.6), Inches(0.45), "2022 개정 교육과정 연계 교과별 탐구 미션 및 교수·학습 활동지 개발, 과학 교구/시약 MSDS 물질안전보건자료 텍스트 데이터베이스 정리", Pt(8.5), False, RGBColor(20, 20, 20))

    add_textbox(s4, Inches(0.65), Inches(4.85), Inches(1.5), Inches(0.4), "팀원 B (디자인/검증)", Pt(9.5), True, RGBColor(0, 80, 180))
    add_textbox(s4, Inches(2.2), Inches(4.81), Inches(4.6), Inches(0.45), "공식 3D 인포그래픽 포스터 디자인, 전자칠판 최적화 진한 굴림체 UI/CSS 튜닝, 학교 현장 시연 테스트 및 사용자 피드백 수집", Pt(8.5), False, RGBColor(20, 20, 20))

    # 13. 완성 기준
    text_p4_4 = (
        "1. 완전 독립형 배포: 인터넷이나 추가 설치 없이 학교 PC에서 NFC에듀태그.exe 더블클릭만으로 완벽 구동.\n"
        "2. NFC 올인원 다중 모드: 출석·감정 체크, 과학실 패스포트 스탬프 발급, 교구 대여/반납이 1초 내에 오류 없이 처리.\n"
        "3. 지능형 과학실 융합: kfcman.link/scibit micro:bit 센서 수집 및 kfcman.link/scibot 햄스터 로봇 탐사가 연동.\n"
        "4. 초거대 AI 및 보안: Upstage Solar AI 피드백이 실시간 생성되고, API 키와 학생 정보가 AES-256 로컬 암호화 보관.\n"
        "5. 데이터 관리: 모든 활동 기록이 엑셀(CSV)로 즉시 다운로드되고 전체 백업/복원이 가능."
    )
    add_textbox(s4, Inches(0.65), Inches(5.95), Inches(6.2), Inches(1.15), text_p4_4, Pt(9), False, RGBColor(20, 20, 20))

    prs.save(out_path1)
    prs.save(out_path2)
    prs.save(out_path3)
    print("성공: 해커톤 팀 활동 기획서 PPTX가 생성되었습니다!")
    print("저장 경로 1:", out_path1)
    print("저장 경로 2:", out_path2)
    print("저장 경로 3:", out_path3)

if __name__ == "__main__":
    fill_hackathon_pptx()
