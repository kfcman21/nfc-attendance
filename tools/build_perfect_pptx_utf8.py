# -*- coding: utf-8 -*-
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_perfect_hackathon_pptx():
    prs = pptx.Presentation()
    # A4 ?몃줈 洹쒓꺽 (210mm x 297mm = 8.27in x 11.69in)
    prs.slide_width = Inches(8.27)
    prs.slide_height = Inches(11.69)
    blank_layout = prs.slide_layouts[6]

    # 怨듯넻 ?됱긽 ?붾젅??    C_DARK = RGBColor(40, 48, 60)        # 硫붿씤 ?ㅽ겕 洹몃젅??    C_YELLOW = RGBColor(254, 209, 86)    # 1/4踰??몃옉
    C_TEAL = RGBColor(85, 214, 190)      # 2踰?泥?줉
    C_PINK = RGBColor(252, 141, 158)     # 3踰??묓겕
    C_BLUE = RGBColor(0, 100, 224)       # ?ъ씤??肄붾컻??釉붾（
    C_BORDER = RGBColor(205, 212, 220)   # 諛뺤뒪 ?멸낸??    C_BOX_BG = RGBColor(255, 255, 255)   # 諛뺤뒪 諛곌꼍
    C_TEXT = RGBColor(25, 30, 36)        # 蹂몃Ц ?띿뒪??    C_MUTED = RGBColor(110, 120, 135)    # 蹂댁“ ?띿뒪??    C_LIGHT_GRAY = RGBColor(245, 247, 250)

    def draw_base_template(slide, page_num, section_title, badge_color):
        # 1. ?곷떒 ?ㅽ겕 ?ㅻ뜑 諛대뱶
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(8.27), Inches(0.55))
        header.fill.solid()
        header.fill.fore_color.rgb = C_DARK
        header.line.color.rgb = C_DARK
        tf = header.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = "?댁빱??? ?쒕룞 湲고쉷??
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # 2. ?二쇱젣 踰덊샇 諛곗?
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(0.8), Inches(0.55), Inches(0.55))
        badge.fill.solid()
        badge.fill.fore_color.rgb = badge_color
        badge.line.fill.background()
        tf_b = badge.text_frame
        tf_b.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_b = tf_b.paragraphs[0]
        p_b.text = str(page_num)
        p_b.font.name = "Malgun Gothic"
        p_b.font.size = Pt(20)
        p_b.font.bold = True
        p_b.font.color.rgb = C_DARK
        p_b.alignment = PP_ALIGN.CENTER

        # 3. ?二쇱젣 ??댄?
        tx_title = slide.shapes.add_textbox(Inches(1.3), Inches(0.8), Inches(4.0), Inches(0.55))
        tf_t = tx_title.text_frame
        tf_t.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_t = tf_t.paragraphs[0]
        p_t.text = section_title
        p_t.font.name = "Malgun Gothic"
        p_t.font.size = Pt(20)
        p_t.font.bold = True
        p_t.font.color.rgb = C_DARK

        # 4. ?섎떒 ?명꽣
        footer_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(11.15), Inches(6.97), Inches(0.01))
        footer_line.fill.solid()
        footer_line.fill.fore_color.rgb = RGBColor(220, 225, 230)
        footer_line.line.color.rgb = RGBColor(220, 225, 230)

        tx_foot_l = slide.shapes.add_textbox(Inches(0.65), Inches(11.2), Inches(5.0), Inches(0.35))
        tf_fl = tx_foot_l.text_frame
        p_fl = tf_fl.paragraphs[0]
        p_fl.text = "2026 湲곗뾽 ?곌퀎 ?뺣낫援먯썝 ??웾媛뺥솕 ?꾨줈?앺듃 쨌 ?댁빱??? 湲고쉷??
        p_fl.font.name = "Malgun Gothic"
        p_fl.font.size = Pt(9.5)
        p_fl.font.color.rgb = C_MUTED

        tx_foot_r = slide.shapes.add_textbox(Inches(6.62), Inches(11.2), Inches(1.0), Inches(0.35))
        tf_fr = tx_foot_r.text_frame
        p_fr = tf_fr.paragraphs[0]
        p_fr.text = f"{page_num} / 4"
        p_fr.font.name = "Malgun Gothic"
        p_fr.font.size = Pt(9.5)
        p_fr.font.bold = True
        p_fr.font.color.rgb = C_MUTED
        p_fr.alignment = PP_ALIGN.RIGHT

    def draw_info_bar(slide, is_first_page=False):
        top_y = Inches(1.48)
        height = Inches(0.7) if is_first_page else Inches(0.4)
        
        info_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), top_y, Inches(6.97), height)
        info_box.fill.solid()
        info_box.fill.fore_color.rgb = C_LIGHT_GRAY
        info_box.line.color.rgb = C_BORDER
        info_box.line.width = Pt(1)
        
        tf = info_box.text_frame
        tf.margin_left = Inches(0.15)
        tf.margin_top = Inches(0.08)
        tf.word_wrap = True
        
        if is_first_page:
            p1 = tf.paragraphs[0]
            p1.text = "遺꾨컲: 珥덈벑 ?뺣낫 ?듯빀 遺꾨컲         紐⑤몺紐? NFC ?먮??쒓렇 ? (Smart EduLab)"
            p1.font.name = "Malgun Gothic"
            p1.font.size = Pt(10.5)
            p1.font.bold = True
            p1.font.color.rgb = C_DARK
            
            p2 = tf.add_paragraph()
            p2.text = "??? 諛뺤갔洹????媛쒕컻 珥앷큵), ???A(援먯쑁怨쇱젙/肄섑뀗痢?, ???B(UI쨌UX ?붿옄???꾩옣寃利?"
            p2.font.name = "Malgun Gothic"
            p2.font.size = Pt(9.5)
            p2.font.color.rgb = C_BLUE
            p2.space_before = Pt(4)
        else:
            p1 = tf.paragraphs[0]
            p1.text = "紐⑤몺紐? NFC ?먮??쒓렇 ? (Smart EduLab)       |       ??? 諛뺤갔洹?
            p1.font.name = "Malgun Gothic"
            p1.font.size = Pt(10)
            p1.font.bold = True
            p1.font.color.rgb = C_DARK

    def draw_section_header(slide, top_y, badge_text, sub_text):
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), top_y, Inches(1.4), Inches(0.32))
        badge.fill.solid()
        badge.fill.fore_color.rgb = C_DARK
        badge.line.fill.background()
        tf_b = badge.text_frame
        tf_b.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_b = tf_b.paragraphs[0]
        p_b.text = badge_text
        p_b.font.name = "Malgun Gothic"
        p_b.font.size = Pt(10.5)
        p_b.font.bold = True
        p_b.font.color.rgb = RGBColor(255, 255, 255)
        p_b.alignment = PP_ALIGN.CENTER

        tx_sub = slide.shapes.add_textbox(Inches(2.15), top_y - Inches(0.02), Inches(5.4), Inches(0.35))
        tf_s = tx_sub.text_frame
        tf_s.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_s = tf_s.paragraphs[0]
        p_s.text = sub_text
        p_s.font.name = "Malgun Gothic"
        p_s.font.size = Pt(10)
        p_s.font.italic = True
        p_s.font.color.rgb = C_MUTED

    def draw_content_box(slide, top_y, height, text_list, font_size=Pt(9.5), line_space=Pt(4)):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), top_y, Inches(6.97), height)
        box.fill.solid()
        box.fill.fore_color.rgb = C_BOX_BG
        box.line.color.rgb = C_BORDER
        box.line.width = Pt(1)
        
        tf = box.text_frame
        tf.margin_left = Inches(0.18)
        tf.margin_right = Inches(0.18)
        tf.margin_top = Inches(0.12)
        tf.margin_bottom = Inches(0.12)
        tf.word_wrap = True
        
        for i, item in enumerate(text_list):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.font.name = "Malgun Gothic"
            p.font.size = font_size
            p.font.color.rgb = C_TEXT
            if i > 0:
                p.space_before = line_space
        return box


    # =========================================================================
    # SLIDE 1: 臾몄젣 李얘린
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    draw_base_template(s1, 1, "臾몄젣 李얘린", C_YELLOW)
    draw_info_bar(s1, is_first_page=True)

    # 1. 臾몄젣 ?뺤쓽
    draw_section_header(s1, Inches(2.32), "1. 臾몄젣 ?뺤쓽", "?꾧뎄???대뼡 ?섏씤?ъ씤?몄씤吏")
    draw_content_box(s1, Inches(2.70), Inches(2.45), [
        "??援먯궗??怨쇱쨷???숆툒 諛?怨쇳븰??愿由??됱젙???뚮え:",
        "  - 留ㅼ씪 ?꾩묠 ?숈깮 異쒓껐 ?섎룞 泥댄겕, 吏媛곸깮 ?뚯븙, ?숈깮?ㅼ쓽 ?뺤꽌쨌留덉쓬 ?곹깭瑜??쇱씪???뚯븙?섍린 ?대젮?.",
        "??吏?ν삎 怨쇳븰???먭뎄 ?대젰 諛??덉쟾 愿由?遺??",
        "  - 怨쇳븰???ㅽ뀒?댁뀡(MBL, ?꾨?寃???蹂??숈깮 ?먭뎄 ?대젰 愿由ъ쓽 ?대젮? 諛?怨좉? 援먭뎄쨌?뷀븰 ?쒖빟(MSDS) ???諛섎궔怨??ㅽ뿕 ?덉쟾 ?섏튃 ?덈궡???뺤떇??",
        "??遺꾩젅???숆툒 愿由??꾧뎄??鍮꾪슚??",
        "  - 異쒖꽍遺, ?낆꽌 ?듭옣, 泥댁쑁(PAPS) 痢≪젙, ?숆툒 ?뷀룓媛 ?쒓컖媛?遺꾨━?섏뼱 援먯궗???섏뾽 以鍮?諛?愿由??쇰줈??媛以?",
        "??援먯쑁??AI??臾대퉬?먯쟻 ?꾩쟾 ?먮룞???꾪뿕:",
        "  - AI媛 ?앹꽦???쇰뱶諛깆씠 援먯궗??寃???놁씠 ?숈깮?먭쾶 洹몃?濡??꾨떖?섏뼱, 援먯궗??二쇱껜?깃낵 ?숈깮??鍮꾪뙋???ш퀬媛 ?④퍡 ?쏀솕??"
    ], font_size=Pt(9.0), line_space=Pt(2))

    # 2. 吏湲?諛⑹떇
    draw_section_header(s1, Inches(5.30), "2. 吏湲?諛⑹떇", "吏湲덉? ?대뼸寃??섍퀬 ?덈뒗吏 쨌 遺덊렪????)
    draw_content_box(s1, Inches(5.68), Inches(2.45), [
        "??異쒓껐 諛?媛먯젙: 醫낆씠 異쒖꽍遺???섏씠??NEIS) ?섎룞 ?낅젰???섏〈?섏뿬 ?숈깮???뺤꽌 蹂?붾굹 ?뺤떆 ?깃탳 異붿씠瑜?吏곴??곸쑝濡??뚯븙?섍린 ?대젮?.",
        "??怨쇳븰???먭뎄: 醫낆씠 ?쒕룞吏 ?꾩옣 ?좎씤 諛??섍린 ?????μ뿉 ?섏〈?섏뿬 援먭뎄 遺꾩떎 ?꾪뿕???믨퀬 ?먭뎄 ?곗씠???꾩쟻 遺덇???",
        "??泥댁쑁 諛??낆꽌: 珥덉떆怨꾨? ?ㅺ퀬 ?섎룞?쇰줈 ?뷀????⑺??꾩쓣 痢≪젙?섍굅??醫낆씠???곸뼱 泥닿퀎?곸씤 ?깆옣 遺꾩꽍???대젮?.",
        "??遺덊렪???? ?대씪?곕뱶 ?뱀빋? 援먯떎 ?명꽣??遺덉븞?????묒냽???딄린硫? ?숈깮 媛쒖씤?뺣낫(?대쫫, ?ъ쭊) ?좎텧 ?곕젮媛 ?곸〈??"
    ], font_size=Pt(9.3), line_space=Pt(3))

    # 3. ?ъ슜??    draw_section_header(s1, Inches(8.28), "3. ?ъ슜??, "?꾧? 쨌 ?대뼡 ?곹솴?먯꽌 ?곕뒗吏")
    draw_content_box(s1, Inches(8.66), Inches(2.25), [
        "??二??ъ슜??1 (珥댟룹쨷??援먯궗):",
        "  - ?꾩묠 ?깃탳 留욎씠(異쒖꽍/媛먯젙 AI 耳??, 吏?ν삎 怨쇳븰???ㅽ뀒?댁뀡 ?먭뎄 ?섏뾽 諛?援먭뎄/?쒖빟 ?덉쟾 ?듭젣, 泥댁쑁 PAPS 痢≪젙, ?숆툒 臾멸퀬 諛?1??1??븷 ?숆툒 ?뷀룓(吏援??ъ씤??瑜??듯빀 ?뺤궛????",
        "??二??ъ슜??2 (珥덈벑?숈깮):",
        "  - ?깃탳 ?? 怨쇳븰???먭뎄 肄붾꼫 ?대룞 ?? 援먭뎄 ?????蹂몄씤??NFC ?숈깮 移대뱶瑜?由щ뜑湲곗뿉 '?? ?쒓렇????"
    ], font_size=Pt(9.3), line_space=Pt(3))


    # =========================================================================
    # SLIDE 2: 湲곕뒫怨??먮즺
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    draw_base_template(s2, 2, "湲곕뒫怨??먮즺", C_TEAL)
    draw_info_bar(s2, is_first_page=False)

    # 4. ?꾩슂 湲곕뒫
    draw_section_header(s2, Inches(2.00), "4. ?꾩슂 湲곕뒫", "?듭떖 湲곕뒫 ?꾩＜ 쨌 ??以꾩뿉 ?섎굹")
    draw_content_box(s2, Inches(2.38), Inches(4.00), [
        "?? NFC ?먰꽣移??뺤떆 異쒖꽍 諛??먯닔??: 留덇컧?쒓컖 ?댁쟾 ?쒓렇 ??留뚯젏 遺??諛??ㅼ떆媛???쒕낫???쒓컖??,
        "?? Upstage Solar AI 媛먯젙 異쒖꽍遺 : 異쒖꽍 ??湲곕텇 ?좏깮 ??珥덇굅? AI??留욎땄??湲띿젙 ?묒썝 硫붿떆吏 ?ㅼ떆媛??앹꽦",
        "?? ?ㅻ쭏??怨쇳븰 ?먭뎄 ?⑥뒪?ы듃 : 怨쇳븰???ㅽ뀒?댁뀡蹂?NFC ?쒓퉭 誘몄뀡 ?꾩＜ ?ㅽ꺃???좊━ ?먮룞??,
        "?? ?ㅻ쭏??援먭뎄 & ?쒖빟 ?덉쟾 吏?댁씠 : NFC 援먭뎄 ?먰꽣移????諛섎궔 諛?MSDS ?ㅽ뿕 ?덉쟾 ?섏튃 ?ㅼ떆媛??앹뾽 ?덈궡",
        "?? SciBit 留덉씠?щ줈鍮꾪듃 MBL ?곌퀎 : micro:bit v2 ?쇱꽌(?⑤룄쨌議곕룄쨌?뚯쓬쨌媛?띾룄) ?ㅼ떆媛??섏쭛 諛?湲곕줉 (kfcman.link/scibit)",
        "?? SciBot ?꾩뒪???쇱?而?AI ?먯궗 ?곌퀎 : ?꾩뒪??濡쒕큸 ?⑥떎 ?쒖같 & ?먯쑉二쇳뻾 ?먯궗 誘몄뀡 ?ㅽ꺃??諛쒓툒 (kfcman.link/scibot)",
        "?? ?ㅻ쭏??PAPS ?뷀???& ?쒗궥 : 諛섑솚???쒓퉭?쇰줈 ?뺣났?щ━湲??먮룞 湲곕줉 諛?Solar AI 泥대젰 肄붿묶 由ы룷??,
        "?? ?먮컮???꾩꽌愿 & 吏援??ъ씤???곸젏 : ?숆툒臾멸퀬 ?먰꽣移????諛?1??1??븷 蹂댁긽 ?숆툒 ?뷀룓 寃쎌젣 ?ㅼ뒿",
        "?? ?멸컙 ?묒뾽??HITL) AI 援먯쑁 ?ㅽ뒠?붿삤 : AI-SPARC ?섏뾽 ?ㅺ퀎湲?+ 援먯궗 Level 3 寃?졖룹듅???쇳꽣 + AI 媛먯궗 異붿쟻 ?쒕룞吏 ?쒖옉湲?,
        "?? 怨듦났?곗씠??& 100% 濡쒖뺄 蹂댁븞 : 湲곗긽泥?룹뿉?댁퐫由ъ븘쨌NEIS ?ㅼ떆媛??곕룞 諛??꾩쟾 ?ㅽ봽?쇱씤 SQLite ?뷀샇??蹂닿?"
    ], font_size=Pt(8.8), line_space=Pt(3))

    # 5. ?낅젰 ?먮즺
    draw_section_header(s2, Inches(6.52), "5. ?낅젰 ?먮즺", "?ｌ쓣 ?먮즺 쨌 ?곗씠??쨌 ?덉떆 臾몄꽌")
    draw_content_box(s2, Inches(6.90), Inches(1.80), [
        "???숈깮 諛?援먭뎄 ?곗씠?? ?숈깮 踰덊샇, 移대뱶 UID, 怨쇳븰???ㅽ뀒?댁뀡 ?뺣낫, 援먭뎄/?쒖빟 蹂닿? ?꾩튂 諛?MSDS ?덉쟾 ?섏튃.",
        "???쇱꽌 諛?怨듦났?곗씠?? micro:bit MBL ?쇱꽌媛? 湲곗긽泥??④린?덈낫(湲곗삩/?좎뵪), ?먯뼱肄붾━??誘몄꽭癒쇱?(PM10/PM2.5), NEIS 湲됱떇.",
        "??蹂댁븞 API ?? Upstage Solar API Key (AES-256-GCM ?뷀샇??蹂닿?)."
    ], font_size=Pt(9.3), line_space=Pt(3))

    # 6. 留뚮뱾 諛⑹떇
    draw_section_header(s2, Inches(8.85), "6. 留뚮뱾 諛⑹떇", "?? ?먮뒗 ?쇳빀")
    draw_content_box(s2, Inches(9.23), Inches(0.65), [
        "???먯씠?꾪듃          ??諛붿씠釉뚯퐫??         ???쇳빀 (Hybrid) ??[Electron/Node 諛붿씠釉뚯퐫??+ Upstage Solar AI ?먯씠?꾪듃]"
    ], font_size=Pt(9.5))

    # 7. ?ъ슜 ?꾧뎄
    draw_section_header(s2, Inches(10.02), "7. ?ъ슜 ?꾧뎄", "???꾧뎄 ?대쫫")
    draw_content_box(s2, Inches(10.40), Inches(0.60), [
        "Node.js, Express, Electron, electron-builder, SQLite3, SerialPort, HTML5/CSS3(吏꾪븳 援대┝泥?, Upstage Solar LLM API"
    ], font_size=Pt(9.5))


    # =========================================================================
    # SLIDE 3: ?숈옉 ?먮쫫
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    draw_base_template(s3, 3, "?숈옉 ?먮쫫", C_PINK)
    draw_info_bar(s3, is_first_page=False)

    draw_section_header(s3, Inches(2.00), "8. ?숈옉 ?먮쫫", "?④퀎蹂??댁슜 쨌 湲 ?먮뒗 洹몃┝")

    # 3??而щ읆 諛뺤뒪 ?앹꽦
    col_w = Inches(2.05)
    col_h = Inches(4.30)
    col_y = Inches(2.40)

    # 而щ읆 1: ?낅젰
    b1 = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), col_y, col_w, col_h)
    b1.fill.solid()
    b1.fill.fore_color.rgb = RGBColor(255, 250, 235)
    b1.line.color.rgb = C_YELLOW
    b1.line.width = Pt(1.5)
    tf1 = b1.text_frame
    tf1.margin_left = Inches(0.12)
    tf1.margin_top = Inches(0.12)
    tf1.word_wrap = True
    p1_h = tf1.paragraphs[0]
    p1_h.text = "1. ????(Input)"
    p1_h.font.name = "Malgun Gothic"
    p1_h.font.size = Pt(11)
    p1_h.font.bold = True
    p1_h.font.color.rgb = RGBColor(160, 110, 0)
    p1_h.alignment = PP_ALIGN.CENTER
    
    items1 = [
        "",
        "??NFC 移대뱶 ?쒓렇",
        "  - ?숈깮 移대뱶 (?숈깮 ?앸퀎)",
        "  - ?ㅽ뀒?댁뀡 移대뱶 (?꾩＜)",
        "  - 援먭뎄 移대뱶 (???諛섎궔)",
        "  - 梨?移대뱶 (?꾩꽌 ???",
        "",
        "??micro:bit MBL ?쇱꽌",
        "  - ?⑤룄, 議곕룄, ?뚯쓬, 媛?띾룄",
        "",
        "???꾩뒪??濡쒕큸 誘몄뀡",
        "  - ?⑥떎 ?쒖같 / ?먯쑉二쇳뻾"
    ]
    for it in items1:
        p = tf1.add_paragraph()
        p.text = it
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(8.8)
        p.font.color.rgb = C_TEXT

    # ?붿궡??1
    arr1 = s3.shapes.add_textbox(Inches(2.75), Inches(4.2), Inches(0.35), Inches(0.5))
    arr1.text_frame.paragraphs[0].text = "??
    arr1.text_frame.paragraphs[0].font.size = Pt(18)
    arr1.text_frame.paragraphs[0].font.bold = True
    arr1.text_frame.paragraphs[0].font.color.rgb = C_MUTED

    # 而щ읆 2: 泥섎━
    b2 = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.11), col_y, col_w, col_h)
    b2.fill.solid()
    b2.fill.fore_color.rgb = RGBColor(235, 252, 248)
    b2.line.color.rgb = C_TEAL
    b2.line.width = Pt(1.5)
    tf2 = b2.text_frame
    tf2.margin_left = Inches(0.12)
    tf2.margin_top = Inches(0.12)
    tf2.word_wrap = True
    p2_h = tf2.paragraphs[0]
    p2_h.text = "2. 泥?由?(Process)"
    p2_h.font.name = "Malgun Gothic"
    p2_h.font.size = Pt(11)
    p2_h.font.bold = True
    p2_h.font.color.rgb = RGBColor(0, 130, 110)
    p2_h.alignment = PP_ALIGN.CENTER
    
    items2 = [
        "",
        "??由щ뜑湲?紐⑤뱶蹂?遺꾧린",
        "  - 異쒖꽍/怨쇳븰???꾩꽌/泥댁쑁",
        "",
        "??濡쒖뺄 SQLite ?몃옖??뀡",
        "  - 異쒓껐, ?ㅽ꺃?? ???湲곕줉",
        "",
        "??Upstage Solar AI 遺꾩꽍",
        "  - 媛먯젙 泥댄겕???묒썝 ?앹꽦",
        "  - ?뷀???留욎땄 肄붿묶",
        "",
        "??HITL Level 3 援먯궗 寃??,
        "  - AI 珥덉븞 鍮꾪뙋쨌?섏젙쨌?뱀씤",
        "",
        "???ㅼ떆媛?SSE 釉뚮줈?쒖틦?ㅽ듃"
    ]
    for it in items2:
        p = tf2.add_paragraph()
        p.text = it
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(8.8)
        p.font.color.rgb = C_TEXT

    # ?붿궡??2
    arr2 = s3.shapes.add_textbox(Inches(5.21), Inches(4.2), Inches(0.35), Inches(0.5))
    arr2.text_frame.paragraphs[0].text = "??
    arr2.text_frame.paragraphs[0].font.size = Pt(18)
    arr2.text_frame.paragraphs[0].font.bold = True
    arr2.text_frame.paragraphs[0].font.color.rgb = C_MUTED

    # 而щ읆 3: 寃곌낵
    b3 = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.57), col_y, col_w, col_h)
    b3.fill.solid()
    b3.fill.fore_color.rgb = RGBColor(255, 242, 245)
    b3.line.color.rgb = C_PINK
    b3.line.width = Pt(1.5)
    tf3 = b3.text_frame
    tf3.margin_left = Inches(0.12)
    tf3.margin_top = Inches(0.12)
    tf3.word_wrap = True
    p3_h = tf3.paragraphs[0]
    p3_h.text = "3. 寃?怨?(Output)"
    p3_h.font.name = "Malgun Gothic"
    p3_h.font.size = Pt(11)
    p3_h.font.bold = True
    p3_h.font.color.rgb = RGBColor(180, 50, 80)
    p3_h.alignment = PP_ALIGN.CENTER
    
    items3 = [
        "",
        "?????紐⑤땲???ㅼ떆媛??덈궡",
        "  - ?뺤떆 異쒖꽍 ?꾨즺 & ?먯닔",
        "  - Solar AI 留욎땄 寃⑸젮 ?앹뾽",
        "  - ?⑥뒪?ы듃 ?ㅽ꺃??諛쒓툒",
        "  - 援먭뎄 ???& MSDS ?쒖텧",
        "",
        "  - 援먯궗 ?뱀씤 ?쇰뱶諛깅쭔 諛고룷",
        "",
        "??AI 媛먯궗 異붿쟻 ?쒕룞吏 異쒕젰",
        "  - ?먭뎄?믪쭏臾멤넂?쒓퉭?믨탳??,
        "",
        "???먰겢由??묒?(CSV) ???,
        "  - ?섏씠???앺솢湲곕줉遺 諛섏쁺",
        "  - 100% 濡쒖뺄 ?곗씠??蹂댁〈"
    ]
    for it in items3:
        p = tf3.add_paragraph()
        p.text = it
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(8.8)
        p.font.color.rgb = C_TEXT

    # 9. 硫섑넗 吏덈Ц
    draw_section_header(s3, Inches(6.90), "9. 硫섑넗 吏덈Ц", "留됲엺 寃?쨌 硫섑넗?먭쾶 臾쇱뼱蹂?寃?)
    draw_content_box(s3, Inches(7.28), Inches(3.60), [
        "1. ?섎뱶?⑥뼱 ?명솚??諛??덉젙??寃利?",
        "   - ?숆탳 而댄벂?곗떎留덈떎 OS ?섍꼍(Win 10/11)怨??쒕━??COM ?ы듃 踰덊샇媛 ?ㅻⅨ?? ?쒕━???ы듃 ?먮룞 ?먯?(Auto-detect) 諛?USB HID ?ㅻ낫??紐⑤뱶 寃몄슜 援ъ“媛 ?ㅼ젣 ?꾩옣?먯꽌 ?덉젙?곸쑝濡?援щ룞?섎뒗吏 ?먭?諛쏄퀬 ?띠뒿?덈떎.",
        "",
        "2. ?쇱?而?而댄벂??援먭뎄 ?ㅼ씠?됲듃 ?곕룞 ?뺤콉:",
        "   - 吏?ν삎 怨쇳븰?ㅼ쓽 留덉씠?щ줈鍮꾪듃(SciBit) 諛??꾩뒪??濡쒕큸(SciBot) ?곗씠?곕? 釉뚮씪?곗???Web Serial/Web Bluetooth API濡??ㅼ씠?됲듃 ?곕룞???뚯쓽 蹂댁븞 媛?대뱶?쇱씤 諛?紐⑤쾾 ?щ???????먮Ц??援ы빀?덈떎.",
        "",
        "3. ?ㅽ봽?쇱씤 AI ?대갚(Fallback) ?꾪궎?띿쿂:",
        "   - ?숆탳 援먯떎 ???명꽣?룹씠 ?쇱떆?곸쑝濡??딄꼈????Upstage Solar AI ?몄텧 ?ㅽ뙣瑜??鍮꾪븯?? ?ㅽ봽?쇱씤 猷곕쿋?댁뒪(Rule-based) 媛먯젙 諛?泥댁쑁 ?쇰뱶諛깆쑝濡?遺?쒕읇寃??꾪솚?섎뒗 Fallback ?ㅺ퀎媛 ?곸젅?쒖? ?됯?諛쏄퀬 ?띠뒿?덈떎."
    ], font_size=Pt(9.0), line_space=Pt(2))


    # =========================================================================
    # SLIDE 4: 紐⑺몴 ?뺥븯湲?    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    draw_base_template(s4, 4, "紐⑺몴 ?뺥븯湲?, C_YELLOW)
    draw_info_bar(s4, is_first_page=False)

    # 10. 寃곌낵臾??대쫫
    draw_section_header(s4, Inches(2.00), "10. 寃곌낵臾??대쫫", "媛移?쨌 ?섏쨷??諛붽퓭????)
    box10 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(2.38), Inches(6.97), Inches(0.65))
    box10.fill.solid()
    box10.fill.fore_color.rgb = RGBColor(238, 246, 255)
    box10.line.color.rgb = C_BLUE
    box10.line.width = Pt(1.5)
    tf10 = box10.text_frame
    tf10.vertical_anchor = MSO_ANCHOR.MIDDLE
    p10 = tf10.paragraphs[0]
    p10.text = "NFC ?먮??쒓렇 (NFC EduTag) - ?ㅻ쭏???숆툒 & 吏?ν삎 怨쇳븰???ъ씤???붾（??
    p10.font.name = "Malgun Gothic"
    p10.font.size = Pt(12)
    p10.font.bold = True
    p10.font.color.rgb = C_BLUE
    p10.alignment = PP_ALIGN.CENTER

    # 11. ??以??뚭컻
    draw_section_header(s4, Inches(3.20), "11. ??以??뚭컻", "留뚮뱾?ㅻ뒗 寃???臾몄옣")
    box11 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(3.58), Inches(6.97), Inches(0.75))
    box11.fill.solid()
    box11.fill.fore_color.rgb = C_LIGHT_GRAY
    box11.line.color.rgb = C_BORDER
    tf11 = box11.text_frame
    tf11.vertical_anchor = MSO_ANCHOR.MIDDLE
    p11 = tf11.paragraphs[0]
    p11.text = '"移대뱶 ???μ쑝濡??숈깮 異쒓껐쨌媛먯젙 耳?대???吏?ν삎 怨쇳븰??MBL ?먭뎄쨌濡쒕큸 ?쇱?而?AI源뚯? ?먰꽣移섎줈 ?닿껐?섎뒗 誘몃옒???ㅻ쭏???ㅼ엥 ?ъ씤???붾（??'
    p11.font.name = "Malgun Gothic"
    p11.font.size = Pt(10)
    p11.font.bold = True
    p11.font.color.rgb = C_DARK
    p11.alignment = PP_ALIGN.CENTER

    # 12. ??븷 遺꾨떞
    draw_section_header(s4, Inches(4.50), "12. ??븷 遺꾨떞", "??먮퀎 ?대떦")
    
    # ??븷 遺꾨떞 ???앹꽦
    tbl_shape = s4.shapes.add_table(4, 2, Inches(0.65), Inches(4.88), Inches(6.97), Inches(2.40))
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(1.8)
    tbl.columns[1].width = Inches(5.17)

    headers = ["??由?, "?대떦 ?낅Т 諛???븷"]
    for c_idx, h_text in enumerate(headers):
        cell = tbl.cell(0, c_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_DARK
        p = cell.text_frame.paragraphs[0]
        p.text = h_text
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    rows_data = [
        ("諛뺤갔洹?n(???/ 媛쒕컻 珥앷큵)", "???꾩껜 ?쒖뒪???꾪궎?띿쿂 ?ㅺ퀎, Electron ?곗뒪?ы넲 ??諛?Node/SQLite DB 援ъ텞\n??吏?ν삎 怨쇳븰??SciBit/SciBot) ?곌퀎 諛?Upstage Solar AI ?붿쭊 ?듯빀 媛쒕컻\n??2022 媛쒖젙 援먯쑁怨쇱젙 ?곌퀎 湲고쉷 諛?理쒖쥌 ?낅┰???⑦궎吏?鍮뚮뱶"),
        ("???A\n(援먯쑁怨쇱젙 / 肄섑뀗痢?", "??2022 媛쒖젙 援먯쑁怨쇱젙 ?곌퀎 援먭낵蹂??먭뎄 誘몄뀡 諛?援먯닔쨌?숈뒿 ?쒕룞吏 媛쒕컻\n??HITL AI-SPARC ?섏뾽 ?ㅺ퀎??諛?4?④퀎 AI 媛먯궗 異붿쟻 ?숈깮 ?쒕룞吏 吏묓븘\n??怨쇳븰 援먭뎄 諛??뷀븰 ?쒖빟 MSDS 臾쇱쭏?덉쟾蹂닿굔?먮즺 ?띿뒪??DB ?뺣━"),
        ("???B\n(UI쨌UX ?붿옄??/ 寃利?", "??怨듭떇 3D ?명룷洹몃옒???ъ뒪???붿옄??諛??꾩옄移좏뙋 理쒖쟻??吏꾪븳 援대┝泥?UI ?쒕떇\n???숆탳 ?꾩옣 ?쒖뿰 ?뚯뒪??諛?湲곕뒫蹂??몃윭釉붿뒋???ъ슜???쇰뱶諛??섏쭛")
    ]

    for r_idx, (name, role) in enumerate(rows_data, start=1):
        cell_n = tbl.cell(r_idx, 0)
        cell_n.fill.solid()
        cell_n.fill.fore_color.rgb = C_LIGHT_GRAY if r_idx % 2 == 1 else RGBColor(255, 255, 255)
        p_n = cell_n.text_frame.paragraphs[0]
        p_n.text = name
        p_n.font.name = "Malgun Gothic"
        p_n.font.size = Pt(9)
        p_n.font.bold = True
        p_n.font.color.rgb = C_BLUE
        p_n.alignment = PP_ALIGN.CENTER

        cell_r = tbl.cell(r_idx, 1)
        cell_r.fill.solid()
        cell_r.fill.fore_color.rgb = C_LIGHT_GRAY if r_idx % 2 == 1 else RGBColor(255, 255, 255)
        cell_r.text_frame.margin_left = Inches(0.1)
        cell_r.text_frame.margin_right = Inches(0.1)
        p_r = cell_r.text_frame.paragraphs[0]
        p_r.text = role
        p_r.font.name = "Malgun Gothic"
        p_r.font.size = Pt(8.5)
        p_r.font.color.rgb = C_TEXT

    # 13. ?꾩꽦 湲곗?
    draw_section_header(s4, Inches(7.45), "13. ?꾩꽦 湲곗?", "?닿쾬???섎㈃ ?꾩꽦")
    draw_content_box(s4, Inches(7.83), Inches(3.10), [
        "1. [?꾩닔] ?꾩쟾 ?낅┰??諛고룷: ?명꽣?룹씠??異붽? ?ㅼ튂 ?놁씠 ?숆탳 PC?먯꽌 NFC?먮??쒓렇.exe ?붾툝?대┃留뚯쑝濡?100% ?뺤긽 援щ룞.",
        "2. [?꾩닔] NFC ?ъ씤???ㅼ쨷 紐⑤뱶: ?숈깮 移대뱶 ?쒓렇 ??異쒖꽍쨌媛먯젙 泥댄겕, 怨쇳븰???⑥뒪?ы듃 ?ㅽ꺃??諛쒓툒, 援먭뎄 ???諛섎궔??1珥??댁뿉 ?ㅻ쪟 ?놁씠 泥섎━.",
        "3. [?꾩닔] 吏?ν삎 怨쇳븰???듯빀: kfcman.link/scibit micro:bit ?쇱꽌 ?섏쭛 諛?kfcman.link/scibot ?꾩뒪??濡쒕큸 ?먯궗 誘몄뀡???⑥뒪?ы듃???뺤긽 ?곕룞.",
        "4. [?꾩닔] 珥덇굅? AI 諛?濡쒖뺄 蹂댁븞: Upstage Solar AI ?쇰뱶諛깆씠 ?ㅼ떆媛??앹꽦?섍퀬, API ?ㅼ? ?숈깮 ?뺣낫媛 AES-256 濡쒖뺄 ?뷀샇?붾줈 ?덉쟾?섍쾶 蹂닿?.",
        "5. [?꾩닔] 援먯궗???곗씠??愿由? 紐⑤뱺 ?쒕룞 湲곕줉???대┃ ??踰덉쑝濡??묒?(CSV)濡?利됱떆 ?ㅼ슫濡쒕뱶?섍퀬 ?꾩껜 諛깆뾽/蹂듭썝???꾨꼍 吏??",
        "6. [?꾩닔] ?멸컙 ?묒뾽??HITL) 寃利? AI ?쇰뱶諛깆씠 援먯궗??Level 3 寃?졖룹듅?몄쓣 嫄곗튇 ?ㅼ뿉留??숈깮?먭쾶 諛고룷?섎ŉ, AI-SPARC ?섏뾽 ?ㅺ퀎?덇낵 4?④퀎 AI 媛먯궗 異붿쟻 ?쒕룞吏媛 ?뺤긽 ?앹꽦."
    ], font_size=Pt(8.8), line_space=Pt(2))

    out_path1 = r"C:\Users\諛뺤갔洹?Desktop\(?꾩꽦蹂??댁빱???_?쒕룞_湲고쉷??NFC?먮??쒓렇.pptx"
    out_path2 = r"C:\Users\諛뺤갔洹?Desktop\Project\nfc\docs\(?꾩꽦蹂??댁빱???_?쒕룞_湲고쉷??NFC?먮??쒓렇.pptx"
    out_path3 = r"C:\Users\諛뺤갔洹?Desktop\Project\nfc\(?꾩꽦蹂??댁빱???_?쒕룞_湲고쉷??NFC?먮??쒓렇.pptx"

    prs.save(out_path1)
    prs.save(out_path2)
    prs.save(out_path3)
    print("?깃났: 湲??留욎땄 理쒖쟻???쇱씠 諛섏쁺???댁빱??湲고쉷??PPTX媛 ?앹꽦?섏뿀?듬땲??")

if __name__ == "__main__":
    create_perfect_hackathon_pptx()
